"""
Generate a 5-minute presentation PPT for the Blockchain Crowdfunding DApp.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # dark navy
ACCENT_BLUE = RGBColor(0x0D, 0x6E, 0xFD)    # bright blue
ACCENT_GREEN = RGBColor(0x00, 0xB8, 0x6E)   # green
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)  # orange
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)     # red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)
CARD_BG = RGBColor(0x25, 0x25, 0x40)        # dark card

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLD_W = Inches(13.333)
SLD_H = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    """Add solid background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, radius=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, lines):
    """
    Add a text box with multiple formatted lines.
    lines: list of (text, font_size, color, bold, alignment) tuples
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text, font_size, color, bold, *rest = line_data
        alignment = rest[0] if rest else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide):
    """Add a footer bar."""
    add_rect(slide, Inches(0), Inches(7.1), SLD_W, Inches(0.4), RGBColor(0x10, 0x10, 0x20))
    add_text_box(slide, Inches(0.5), Inches(7.15), Inches(12), Inches(0.3),
                 "区块链众筹系统 — Blockchain Crowdfunding DApp",
                 font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)


def slide_number(slide, num):
    add_text_box(slide, Inches(12.3), Inches(7.15), Inches(0.8), Inches(0.3),
                 str(num), font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BG)

# Decorative shapes
add_rect(slide1, Inches(0), Inches(0), SLD_W, Inches(0.08), ACCENT_BLUE)
add_rect(slide1, Inches(0), Inches(7.42), SLD_W, Inches(0.08), ACCENT_GREEN)

# Large circle decoration (top-right)
circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x0D, 0x6E, 0xFD)
circle.line.fill.background()
circle.fill.fore_color.brightness = 0.85

# Title
add_text_box(slide1, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             "区块链众筹系统", font_size=56, color=WHITE, bold=True)
add_accent_bar(slide1, Inches(1), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_BLUE)
add_text_box(slide1, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "Crowdfunding DApp — 基于以太坊智能合约的去中心化众筹平台",
             font_size=24, color=LIGHT_GRAY)

# Info line
add_rich_text(slide1, Inches(1), Inches(5.0), Inches(11), Inches(1.5), [
    ("技术栈: Solidity · Hardhat · ethers.js · MetaMask · Bootstrap 5", 16, ACCENT_BLUE, False),
    ("", 10, WHITE, False),
    ("项目展示 · 2026年6月", 14, DARK_GRAY, False),
])

slide_number(slide1, 1)

# ============================================================
# SLIDE 2: Project Overview
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, DARK_BG)
add_bottom_bar(slide2)
add_rect(slide2, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📋 项目概述", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide2, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# What problem it solves
add_rich_text(slide2, Inches(0.8), Inches(1.5), Inches(7.5), Inches(2.5), [
    ("🎯 解决什么问题？", 22, ACCENT_BLUE, True),
    ("", 8, WHITE, False),
    ("传统众筹平台依赖中心化机构，存在透明度低、手续费高、资金托管风险大等问题。", 16, LIGHT_GRAY, False),
    ("本系统利用区块链技术，实现了完全去中心化的众筹流程：", 16, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("  ✅  智能合约自动执行，无需第三方托管资金", 15, WHITE, False),
    ("  ✅  所有交易记录上链，公开透明不可篡改", 15, WHITE, False),
    ("  ✅  捐赠者可随时查看资金流向和项目进度", 15, WHITE, False),
    ("  ✅  项目失败自动退款，保障捐赠者权益", 15, WHITE, False),
])

# Core features card
card = add_rect(slide2, Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.2), CARD_BG, 0.05)
add_text_box(slide2, Inches(9.1), Inches(1.7), Inches(3.5), Inches(0.5),
             "🔑 核心功能", font_size=20, color=ACCENT_ORANGE, bold=True)

add_rich_text(slide2, Inches(9.1), Inches(2.3), Inches(3.5), Inches(4.0), [
    ("① 创建众筹项目", 15, WHITE, True),
    ("   设定目标金额与截止日期", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("② 捐赠 ETH", 15, WHITE, True),
    ("   支持任意金额的以太坊转账", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("③ 结束项目", 15, WHITE, True),
    ("   到期后自动判断成败", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("④ 提取资金 / 退款", 15, WHITE, True),
    ("   成功→发起人提款", 13, LIGHT_GRAY, False),
    ("   失败→捐赠者退款", 13, LIGHT_GRAY, False),
])

slide_number(slide2, 2)

# ============================================================
# SLIDE 3: Technical Architecture
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, DARK_BG)
add_bottom_bar(slide3)
add_rect(slide3, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🏗️ 技术架构", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide3, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# 3-tier architecture
layers = [
    ("🎨 前端展示层", ACCENT_BLUE,
     ["HTML5 + CSS3 + Bootstrap 5 响应式 UI",
      "ethers.js v6 与区块链交互",
      "MetaMask 钱包连接与网络管理",
      "单页面应用，三大功能面板"]),
    ("🔗 区块链交互层", ACCENT_GREEN,
     ["Hardhat 本地开发节点 (Chain ID: 31337)",
      "ethers.js BrowserProvider + JsonRpcProvider",
      "合约 ABI 接口封装与调用",
      "事件监听与交易状态追踪"]),
    ("📜 智能合约层", ACCENT_ORANGE,
     ["Solidity 0.8.28 编写",
      "完整的项目生命周期管理",
      "安全的资金托管与转账逻辑",
      "可查询的事件日志与状态存储"]),
]

for i, (title, color, items) in enumerate(layers):
    left = Inches(0.8 + i * 4.0)
    top = Inches(1.6)

    # Layer card
    card = add_rect(slide3, left, top, Inches(3.6), Inches(5.2), CARD_BG, 0.05)
    add_accent_bar(slide3, left, top, Inches(3.6), Inches(0.05), color)

    add_text_box(slide3, left + Inches(0.3), top + Inches(0.2), Inches(3.0), Inches(0.5),
                 title, font_size=20, color=color, bold=True)

    for j, item in enumerate(items):
        add_text_box(slide3, left + Inches(0.3), top + Inches(1.1 + j * 0.55),
                     Inches(3.0), Inches(0.5),
                     f"• {item}", font_size=13, color=LIGHT_GRAY)

    # Arrow between layers
    if i < 2:
        arrow = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        left + Inches(3.6), Inches(3.8),
                                        Inches(0.4), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

# Architecture flow description
add_text_box(slide3, Inches(0.8), Inches(6.9), Inches(12), Inches(0.3),
             "用户浏览器 ← MetaMask → ethers.js → Hardhat 本地节点 → 智能合约 → 区块链状态",
             font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(slide3, 3)

# ============================================================
# SLIDE 4: Smart Contract Design
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, DARK_BG)
add_bottom_bar(slide4)
add_rect(slide4, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📜 智能合约设计 — Crowdfunding.sol", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide4, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# Data structures
add_text_box(slide4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "📦 数据结构", font_size=22, color=ACCENT_BLUE, bold=True)

code_style_left = Inches(0.8)
code_top = Inches(2.1)

# Struct: Project
add_rect(slide4, code_style_left, code_top, Inches(5.5), Inches(2.1), CARD_BG, 0.03)
add_rich_text(slide4, code_style_left + Inches(0.2), code_top + Inches(0.1), Inches(5.1), Inches(1.9), [
    ("struct Project {", 13, ACCENT_GREEN, True),
    ("    uint256 id;              // 项目编号", 12, LIGHT_GRAY, False),
    ("    string name;             // 项目名称", 12, LIGHT_GRAY, False),
    ("    string description;      // 项目描述", 12, LIGHT_GRAY, False),
    ("    address payable creator; // 发起人地址", 12, LIGHT_GRAY, False),
    ("    uint256 targetAmount;    // 目标金额 (wei)", 12, LIGHT_GRAY, False),
    ("    uint256 currentAmount;   // 已筹金额 (wei)", 12, LIGHT_GRAY, False),
    ("    uint256 deadline;        // 截止时间戳", 12, LIGHT_GRAY, False),
    ("    ProjectStatus status;    // 项目状态", 12, LIGHT_GRAY, False),
    ("}", 13, ACCENT_GREEN, True),
])

# Struct: Donation
add_rect(slide4, code_style_left, Inches(4.3), Inches(5.5), Inches(1.3), CARD_BG, 0.03)
add_rich_text(slide4, code_style_left + Inches(0.2), Inches(4.4), Inches(5.1), Inches(1.1), [
    ("struct Donation {", 13, ACCENT_GREEN, True),
    ("    address donor;      // 捐赠者地址", 12, LIGHT_GRAY, False),
    ("    uint256 amount;      // 捐赠金额", 12, LIGHT_GRAY, False),
    ("    uint256 timestamp;   // 捐赠时间", 12, LIGHT_GRAY, False),
    ("}", 13, ACCENT_GREEN, True),
])

# State machine
add_text_box(slide4, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "🔄 项目状态机", font_size=22, color=ACCENT_ORANGE, bold=True)

states = [
    ("Ongoing\n进行中", ACCENT_BLUE),
    ("Successful\n已成功", ACCENT_GREEN),
    ("Failed\n已失败", ACCENT_RED),
]

for i, (label, color) in enumerate(states):
    left = Inches(7.0 + i * 2.1)
    top = Inches(2.2)
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(1.5), Inches(1.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    if i < 2:
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        left + Inches(1.5), top + Inches(0.55),
                                        Inches(0.6), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

# Transition conditions
add_rich_text(slide4, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.5), [
    ("状态转换条件：", 16, WHITE, True),
    ("", 6, WHITE, False),
    ("• 创建项目 → Ongoing（自动）", 14, LIGHT_GRAY, False),
    ("• Ongoing → Successful", 14, ACCENT_GREEN, False),
    ("   条件: 截止时间已过 AND 已筹 ≥ 目标", 12, DARK_GRAY, False),
    ("• Ongoing → Failed", 14, ACCENT_RED, False),
    ("   条件: 截止时间已过 AND 已筹 < 目标", 12, DARK_GRAY, False),
    ("• 任何人可调用 endProject() 触发状态变更", 14, LIGHT_GRAY, False),
])

slide_number(slide4, 4)

# ============================================================
# SLIDE 5: Core Business Flow
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, DARK_BG)
add_bottom_bar(slide5)
add_rect(slide5, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "⚡ 核心业务流程", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide5, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# Flow steps - horizontal timeline
steps = [
    ("1️⃣", "创建项目", "发起人设定名称、\n描述、目标金额、\n截止日期", ACCENT_BLUE),
    ("2️⃣", "开始捐赠", "任何用户可向\n项目捐赠ETH,\n金额不限", ACCENT_GREEN),
    ("3️⃣", "截止判定", "到期后任何人\n可触发项目\n结算", ACCENT_ORANGE),
    ("4️⃣", "资金分配", "成功→发起人提款\n失败→捐赠者退款", ACCENT_RED if True else ACCENT_RED),
]

for i, (icon, title, desc, color) in enumerate(steps):
    left = Inches(0.6 + i * 3.15)
    top = Inches(1.6)

    # Step card
    card = add_rect(slide5, left, top, Inches(2.8), Inches(3.8), CARD_BG, 0.05)
    add_accent_bar(slide5, left, Inches(1.6), Inches(2.8), Inches(0.05), color)

    # Step number
    add_text_box(slide5, left + Inches(0.3), Inches(1.8), Inches(2.2), Inches(0.6),
                 icon, font_size=28, color=color, bold=True)

    # Step title
    add_text_box(slide5, left + Inches(0.3), Inches(2.4), Inches(2.2), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True)

    # Step description
    add_text_box(slide5, left + Inches(0.3), Inches(3.0), Inches(2.2), Inches(1.5),
                 desc, font_size=13, color=LIGHT_GRAY)

    # Connector arrow
    if i < 3:
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        left + Inches(2.8), Inches(3.2),
                                        Inches(0.35), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

# Key points
add_text_box(slide5, Inches(0.8), Inches(5.7), Inches(12), Inches(0.5),
             "🔍 关键函数调用流程", font_size=20, color=ACCENT_BLUE, bold=True)

flow_items = [
    ("createProject(name, desc, target, deadline)", "发起项目", ACCENT_BLUE),
    ("donate(projectId) {value: ETH}", "捐赠资金", ACCENT_GREEN),
    ("endProject(projectId)", "结算项目", ACCENT_ORANGE),
    ("withdrawFunds(projectId)", "发起人提款", ACCENT_GREEN),
    ("claimRefund(projectId)", "捐赠者退款", ACCENT_RED),
]

for i, (func, label, color) in enumerate(flow_items):
    left = Inches(0.8 + i * 2.5)
    add_rect(slide5, left, Inches(6.2), Inches(2.2), Inches(0.6), color, 0.1)
    add_text_box(slide5, left + Inches(0.1), Inches(6.22), Inches(2.0), Inches(0.55),
                 f"{label}\n{func}", font_size=9, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

slide_number(slide5, 5)

# ============================================================
# SLIDE 6: Frontend Architecture
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, DARK_BG)
add_bottom_bar(slide6)
add_rect(slide6, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🖥️ 前端架构设计", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide6, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# Tab system
add_text_box(slide6, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "📑 三大功能面板", font_size=22, color=ACCENT_BLUE, bold=True)

tabs = [
    ("📊", "所有项目", "项目卡片列表\n进度条可视化\n操作按钮（捐赠/结束/提款/退款）\n捐赠者列表弹窗", ACCENT_BLUE),
    ("➕", "创建项目", "表单填写\n(名称/描述/金额/截止日期)\n表单验证\n交易确认", ACCENT_GREEN),
    ("👤", "我的信息", "钱包地址与余额\n我的捐赠记录\n我创建的项目\n实时数据刷新", ACCENT_ORANGE),
]

for i, (icon, title, desc, color) in enumerate(tabs):
    left = Inches(0.8 + i * 4.0)
    top = Inches(2.2)
    card = add_rect(slide6, left, top, Inches(3.6), Inches(3.4), CARD_BG, 0.05)
    add_accent_bar(slide6, left, top, Inches(3.6), Inches(0.05), color)
    add_text_box(slide6, left + Inches(0.3), top + Inches(0.2), Inches(3.0), Inches(0.5),
                 f"{icon} {title}", font_size=20, color=color, bold=True)
    add_text_box(slide6, left + Inches(0.3), top + Inches(0.9), Inches(3.0), Inches(2.3),
                 desc, font_size=13, color=LIGHT_GRAY)

# Two-provider architecture
add_text_box(slide6, Inches(0.8), Inches(5.8), Inches(5.8), Inches(0.5),
             "🔌 双 Provider 架构（核心创新）", font_size=22, color=ACCENT_ORANGE, bold=True)

# Left: MetaMask provider
add_rect(slide6, Inches(0.8), Inches(6.4), Inches(3.5), Inches(0.6), ACCENT_BLUE, 0.05)
add_text_box(slide6, Inches(1.0), Inches(6.42), Inches(3.1), Inches(0.55),
             "MetaMask Provider → 读取数据 + 钱包UI",
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide6, Inches(4.4), Inches(6.42), Inches(0.5), Inches(0.55),
             "+", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right: Direct provider
add_rect(slide6, Inches(5.0), Inches(6.4), Inches(4.0), Inches(0.6), ACCENT_GREEN, 0.05)
add_text_box(slide6, Inches(5.2), Inches(6.42), Inches(3.6), Inches(0.55),
             "JsonRpcProvider → 发送交易 (无弹窗)",
             font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Explanation
add_text_box(slide6, Inches(9.5), Inches(6.4), Inches(3.5), Inches(0.6),
             "Hardhat 模拟账户 → 自动签名\n开发体验：零弹窗确认 ✅",
             font_size=11, color=LIGHT_GRAY)

slide_number(slide6, 6)

# ============================================================
# SLIDE 7: Security Design
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, DARK_BG)
add_bottom_bar(slide7)
add_rect(slide7, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🔒 安全设计", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide7, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# Security topics
security_items = [
    ("🛡️ 重入攻击防护 (Re-entrancy Guard)",
     ACCENT_RED,
     ["采用 Checks-Effects-Interactions 经典模式",
      "先更新状态变量 (currentAmount = 0)",
      "再执行外部调用 (call{value: amount})",
      "即使攻击者递归调用，状态已清零，无法重复提款"]),
    ("🔐 访问控制 (Access Control)",
     ACCENT_ORANGE,
     ["modifier projectExists: 确保项目存在",
      "require(msg.sender == p.creator): 仅发起人可提款",
      "require(status == Successful): 成功后才可以提款",
      "require(status == Failed): 失败后才可以退款"]),
    ("✅ 输入验证 (Input Validation)",
     ACCENT_GREEN,
     ["项目名称不能为空",
      "目标金额必须大于 0",
      "截止日期必须在未来",
      "捐赠金额必须大于 0",
      "捐赠时项目必须是 Ongoing 状态"]),
    ("📊 事件系统 (Event Logging)",
     ACCENT_BLUE,
     ["ProjectCreated / DonationReceived",
      "ProjectEnded / FundsWithdrawn / RefundClaimed",
      "所有关键操作均触发索引事件",
      "前端可实时监听区块链状态变更"]),
]

for i, (title, color, items) in enumerate(security_items):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.5 + row * 2.85)

    card = add_rect(slide7, left, top, Inches(6.0), Inches(2.55), CARD_BG, 0.05)
    add_accent_bar(slide7, left, top, Inches(0.06), Inches(2.55), color)

    add_text_box(slide7, left + Inches(0.3), top + Inches(0.15), Inches(5.4), Inches(0.5),
                 title, font_size=18, color=color, bold=True)

    for j, item in enumerate(items):
        add_text_box(slide7, left + Inches(0.3), top + Inches(0.8 + j * 0.4),
                     Inches(5.4), Inches(0.4),
                     f"• {item}", font_size=12, color=LIGHT_GRAY)

# Code snippet at bottom
add_rect(slide7, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.55), RGBColor(0x10, 0x10, 0x20), 0.03)
add_text_box(slide7, Inches(0.7), Inches(6.48), Inches(11.9), Inches(0.5),
             "// Checks-Effects-Interactions — 防重入核心模式\n"
             "uint256 amount = p.currentAmount;  // Checks\n"
             "p.currentAmount = 0;               // Effects (先清零!)\n"
             "(bool success, ) = p.creator.call{value: amount}(\"\");  // Interactions (后转账!)",
             font_size=10, color=ACCENT_GREEN)

slide_number(slide7, 7)

# ============================================================
# SLIDE 8: Development Toolchain
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, DARK_BG)
add_bottom_bar(slide8)
add_rect(slide8, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🛠️ 开发工具链与运行流程", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide8, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

# Tool stack
tools = [
    ("Hardhat", "以太坊开发框架\n编译、部署、测试\n本地测试网络", "⚙️", ACCENT_ORANGE),
    ("Solidity", "智能合约语言\nv0.8.28\nviaIR 编译优化", "📜", ACCENT_BLUE),
    ("ethers.js v6", "区块链交互库\nBrowserProvider\nJsonRpcProvider", "🔗", ACCENT_GREEN),
    ("MetaMask", "浏览器钱包插件\n账户管理\n交易签名", "🦊", ACCENT_ORANGE),
    ("Bootstrap 5", "响应式 UI 框架\n卡片、模态框\nToast 通知", "🎨", ACCENT_BLUE),
]

for i, (name, desc, icon, color) in enumerate(tools):
    left = Inches(0.5 + i * 2.55)
    card = add_rect(slide8, left, Inches(1.5), Inches(2.3), Inches(2.2), CARD_BG, 0.05)
    add_accent_bar(slide8, left, Inches(1.5), Inches(2.3), Inches(0.05), color)
    add_text_box(slide8, left + Inches(0.15), Inches(1.65), Inches(2.0), Inches(0.5),
                 f"{icon} {name}", font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, left + Inches(0.15), Inches(2.2), Inches(2.0), Inches(1.3),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Run flow
add_text_box(slide8, Inches(0.8), Inches(4.0), Inches(12), Inches(0.5),
             "📋 本地运行步骤", font_size=22, color=ACCENT_BLUE, bold=True)

steps_list = [
    ("1", "npm install", "安装依赖", ACCENT_BLUE),
    ("2", "npx hardhat node", "启动本地链 (localhost:8545)", ACCENT_GREEN),
    ("3", "npm run deploy", "部署合约到本地链", ACCENT_ORANGE),
    ("4", "npm run dev", "启动前端 (localhost:3000)", ACCENT_BLUE),
    ("5", "连接 MetaMask", "导入 Hardhat 账户到钱包", ACCENT_GREEN),
]

for i, (num, cmd, desc, color) in enumerate(steps_list):
    left = Inches(0.8 + i * 2.5)
    add_rect(slide8, left, Inches(4.7), Inches(2.2), Inches(1.8), CARD_BG, 0.05)
    add_accent_bar(slide8, left, Inches(4.7), Inches(2.2), Inches(0.05), color)
    add_text_box(slide8, left, Inches(4.85), Inches(2.2), Inches(0.5),
                 f"Step {num}", font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, left + Inches(0.1), Inches(5.3), Inches(2.0), Inches(0.6),
                 cmd, font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, left + Inches(0.1), Inches(5.8), Inches(2.0), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Custom tasks
add_text_box(slide8, Inches(0.8), Inches(6.7), Inches(12), Inches(0.3),
             "💡 自定义 Hardhat 任务:  npx hardhat chainInfo  → 查看链信息  |  npx hardhat accounts  → 查看账户余额",
             font_size=12, color=DARK_GRAY)

slide_number(slide8, 8)

# ============================================================
# SLIDE 9: Key Technical Highlights
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, DARK_BG)
add_bottom_bar(slide9)
add_rect(slide9, Inches(0), Inches(0), SLD_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide9, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🌟 关键技术亮点", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide9, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), ACCENT_BLUE)

highlights = [
    ("🔌 双 Provider 架构",
     "前端同时维护 MetaMask Provider（读数据 + 钱包 UI）和\nJsonRpcProvider（发送交易），利用 Hardhat 的\nhardhat_impersonateAccount 实现零弹窗交易确认，\n极大优化开发调试体验。",
     ACCENT_BLUE),
    ("🔄 智能网络检测与切换",
     "自动检测 MetaMask 当前网络 Chain ID，若不在\nHardhat TestNet (31337) 则自动请求切换；\n若网络未添加则自动通过 wallet_addEthereumChain\n添加，确保用户始终在正确的网络上操作。",
     ACCENT_GREEN),
    ("📡 事件驱动架构",
     "智能合约定义 5 个关键事件，前端通过监听\naccountsChanged 和 chainChanged 事件实时响应\n钱包状态变化，自动刷新数据无需手动操作。",
     ACCENT_ORANGE),
    ("🔐 Checks-Effects-Interactions",
     "提款和退款函数严格遵循 CEI 模式：\n① 读取金额 → ② 清零余额 → ③ 外部转账\n有效防止重入攻击，保障资金安全。",
     ACCENT_RED),
    ("💰 自动资金注入",
     "当检测到用户 MetaMask 地址余额不足时，自动从\nHardhat 预置的富账户转账 1000 ETH，确保\n测试用户有足够资金进行各种操作。",
     ACCENT_GREEN),
    ("📱 响应式 UI + 状态可视化",
     "Bootstrap 5 响应式布局适配手机/平板/桌面；\n进度条实时显示筹资比例；项目卡片通过颜色和\n状态徽章直观展示 Ongoing/Successful/Failed 状态。",
     ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(highlights):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.5 + row * 2.9)

    card = add_rect(slide9, left, top, Inches(3.9), Inches(2.6), CARD_BG, 0.05)
    add_accent_bar(slide9, left, top, Inches(3.9), Inches(0.05), color)

    add_text_box(slide9, left + Inches(0.25), top + Inches(0.2), Inches(3.4), Inches(0.5),
                 title, font_size=17, color=color, bold=True)

    add_text_box(slide9, left + Inches(0.25), top + Inches(0.8), Inches(3.4), Inches(1.7),
                 desc, font_size=11, color=LIGHT_GRAY)

slide_number(slide9, 9)

# ============================================================
# SLIDE 10: Summary & Demo
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, DARK_BG)
add_rect(slide10, Inches(0), Inches(0), SLD_W, Inches(0.08), ACCENT_BLUE)
add_rect(slide10, Inches(0), Inches(7.42), SLD_W, Inches(0.08), ACCENT_GREEN)

# Title
add_text_box(slide10, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "🎯 总结与展望", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide10, Inches(0.8), Inches(1.5), Inches(2.0), Inches(0.05), ACCENT_BLUE)

# Summary points
add_rich_text(slide10, Inches(0.8), Inches(2.0), Inches(7.0), Inches(4.0), [
    ("✅ 项目成果", 24, ACCENT_GREEN, True),
    ("", 10, WHITE, False),
    ("• 完整的去中心化众筹 DApp，涵盖项目全生命周期", 17, WHITE, False),
    ("• Solidity 智能合约 + Hardhat 开发环境 + 前端三件套", 17, WHITE, False),
    ("• 安全可靠：CEI 模式防重入 + 严格访问控制", 17, WHITE, False),
    ("• 开发友好：双 Provider 架构实现零弹窗交易", 17, WHITE, False),
    ("• 用户体验：自动网络切换 + 实时状态更新 + 响应式UI", 17, WHITE, False),
    ("", 12, WHITE, False),
    ("🔮 未来展望", 24, ACCENT_BLUE, True),
    ("", 10, WHITE, False),
    ("• 支持 ERC-20 代币捐赠（多币种众筹）", 15, LIGHT_GRAY, False),
    ("• 引入 DAO 治理机制，捐赠者投票决定资金用途", 15, LIGHT_GRAY, False),
    ("• 部署到以太坊测试网 (Sepolia/Goerli)", 15, LIGHT_GRAY, False),
    ("• 添加 NFT 权益证明（捐赠凭证 NFT）", 15, LIGHT_GRAY, False),
    ("• 集成 The Graph 协议实现高效链上数据索引", 15, LIGHT_GRAY, False),
])

# Tech stack summary on the right
add_rect(slide10, Inches(8.5), Inches(2.0), Inches(4.3), Inches(4.5), CARD_BG, 0.05)
add_text_box(slide10, Inches(8.8), Inches(2.2), Inches(3.7), Inches(0.5),
             "🛠️ 技术栈总结", font_size=20, color=ACCENT_ORANGE, bold=True)

tech_items = [
    ("Smart Contract", "Solidity 0.8.28"),
    ("Dev Framework", "Hardhat 2.22"),
    ("Frontend Lib", "ethers.js v6"),
    ("Wallet", "MetaMask"),
    ("UI Framework", "Bootstrap 5.3"),
    ("CSS", "Custom + Bootstrap Icons"),
    ("Node", "http-server"),
    ("Chain", "Hardhat Local (31337)"),
]

for i, (label, value) in enumerate(tech_items):
    y = Inches(2.9 + i * 0.45)
    add_text_box(slide10, Inches(8.8), y, Inches(1.8), Inches(0.4),
                 label, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide10, Inches(10.6), y, Inches(2.0), Inches(0.4),
                 value, font_size=12, color=WHITE, bold=True)

# Demo QR / Video hint
add_text_box(slide10, Inches(0.8), Inches(6.5), Inches(7.0), Inches(0.5),
             "📹 演示视频: 视频演示.mp4  |  📁 源码: contracts/ + src/ + scripts/",
             font_size=14, color=DARK_GRAY)

# Thank you
add_text_box(slide10, Inches(9.0), Inches(6.8), Inches(3.5), Inches(0.4),
             "谢谢！🙏", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(slide10, 10)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "区块链众筹系统-项目介绍.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
